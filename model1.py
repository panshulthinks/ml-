import os
import google.generativeai as genai
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import requests
from PIL import Image
import io
import json

print(load_dotenv())

class PPTGenerator:
    def __init__(self, api_key=None):
        """Initialize the PPT Generator with Gemini API"""
        self.api_key = api_key or os.getenv('GEMINI_API_KEY')
        if not self.api_key:
            raise ValueError("Gemini API key is required. Set GEMINI_API_KEY environment variable or pass it to the constructor.")

        # Configure Gemini
        genai.configure(api_key=self.api_key)
        self.model = genai.GenerativeModel('gemini-1.5-flash')
        self.model_vision = genai.GenerativeModel('gemini-1.5-flash')

        # Initialize presentation
        self.presentation = Presentation()

    def generate_content_outline(self, topic, num_slides=5):
        """Generate content outline using Gemini"""
        prompt = f"""
        Create a detailed outline for a PowerPoint presentation on "{topic}" with {num_slides} slides.
        Return the response as a JSON array with the following structure:
        [
            {{
                "title": "Slide Title",
                "content": "Main content points as bullet points",
                "slide_type": "title|content|image|conclusion"
            }}
        ]

        Make sure the content is engaging, informative, and well-structured.
        The response must be a valid JSON array.
        """

        try:
            response = self.model.generate_content(prompt)
            if hasattr(response, 'text'):
                content = response.text.strip()
            else:
                content = str(response).strip()

            if isinstance(content, list):
                return content[:num_slides]

            # Handle potential markdown code block
            if "```json" in content:
                content = content.split("```json")[1].split("```")[0].strip()
            elif "```" in content:
                content = content.split("```")[1].strip()

            # Ensure we have a valid JSON string
            if not content.startswith('[') or not content.endswith(']'):
                print("Invalid JSON format received. Using fallback outline.")
                return self._get_fallback_outline(topic, num_slides)

            try:
                return json.loads(content)
            except json.JSONDecodeError as je:
                print(f"JSON parsing error: {je}")
                return self._get_fallback_outline(topic, num_slides)

        except Exception as e:
            print(f"Error generating content: {e}")
            return self._get_fallback_outline(topic, num_slides)

    def _get_fallback_outline(self, topic, num_slides):
        """Fallback outline if Gemini fails"""
        return [
            {
                "title": f"Introduction to {topic}",
                "content": "• Overview of the topic\n• Key objectives\n• What to expect",
                "slide_type": "title"
            },
            {
                "title": "Background Information",
                "content": "• Historical context\n• Current state\n• Important facts",
                "slide_type": "content"
            },
            {
                "title": "Key Concepts",
                "content": "• Main principles\n• Core ideas\n• Essential knowledge",
                "slide_type": "content"
            },
            {
                "title": "Applications and Examples",
                "content": "• Real-world applications\n• Case studies\n• Practical examples",
                "slide_type": "content"
            },
            {
                "title": "Conclusion",
                "content": "• Summary of key points\n• Final thoughts\n• Next steps",
                "slide_type": "conclusion"
            }
        ][:num_slides]

    def generate_image_description(self, slide_content):
        """Generate image description for slides using Gemini"""
        prompt = f"""
        Based on this slide content, suggest a relevant image description that would enhance the presentation:

        {slide_content}

        Return only a brief, descriptive phrase suitable for image search (max 10 words).
        """

        try:
            response = self.model.generate_content(prompt)
            return response.text.strip()
        except Exception as e:
            print(f"Error generating image description: {e}")
            return "professional presentation"

    def download_image(self, query, save_path="temp_image.jpg"):
        """Download a relevant image from Pexels based on the query"""
        try:
            # Pexels API endpoint
            url = "https://api.pexels.com/v1/search"

            # Headers with API key
            headers = {
                'Authorization': ''
            }

            # Search parameters
            params = {
                'query': query,
                'per_page': 1,
                'orientation': 'landscape'
            }

            # Get search results
            response = requests.get(url, headers=headers, params=params)
            response.raise_for_status()

            # Parse response
            data = response.json()
            if not data.get('photos'):
                raise ValueError(f"No images found for query: {query}")

            # Get the first photo's original size URL
            image_url = data['photos'][0]['src']['original']

            # Download the image
            img_response = requests.get(image_url)
            img_response.raise_for_status()

            # Save the image
            with open(save_path, 'wb') as f:
                f.write(img_response.content)

            return save_path

        except Exception as e:
            print(f"Error downloading image: {e}")
            # Fallback to placeholder if there's an error
            try:
                img = Image.new('RGB', (800, 600), color='#4A90E2')
                img.save(save_path)
                return save_path
            except Exception as e:
                print(f"Error creating placeholder image: {e}")
                return None

    def create_title_slide(self, title, subtitle=""):
        """Create a title slide"""
        slide_layout = self.presentation.slide_layouts[0]  # Title slide layout
        slide = self.presentation.slides.add_slide(slide_layout)

        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Set subtitle if provided
        if subtitle:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
            subtitle_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        return slide

    def create_content_slide(self, title, content, include_image=False):
        """Create a content slide with bullet points"""
        slide_layout = self.presentation.slide_layouts[1]  # Title and content layout
        slide = self.presentation.slides.add_slide(slide_layout)

        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.bold = True

        # Set content
        content_shape = slide.placeholders[1]
        content_shape.text = content

        # Format bullet points
        text_frame = content_shape.text_frame
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = RGBColor(51, 51, 51)

        # Add image if requested
        if include_image:
            try:
                image_desc = self.generate_image_description(content)
                image_path = self.download_image(image_desc)
                if image_path and os.path.exists(image_path):
                    slide.shapes.add_picture(image_path, Inches(6), Inches(2), height=Inches(4))
                    # Clean up temporary image
                    os.remove(image_path)
            except Exception as e:
                print(f"Error adding image: {e}")

        return slide

    def create_image_slide(self, title, content, image_query):
        """Create a slide with image and text"""
        slide_layout = self.presentation.slide_layouts[8]  # Blank layout
        slide = self.presentation.slides.add_slide(slide_layout)

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(36)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Add content
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4), Inches(5))
        content_frame = content_box.text_frame
        content_frame.text = content

        # Format content
        for paragraph in content_frame.paragraphs:
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = RGBColor(51, 51, 51)

        # Add image
        try:
            image_path = self.download_image(image_query)
            if image_path and os.path.exists(image_path):
                slide.shapes.add_picture(image_path, Inches(5), Inches(2), height=Inches(5))
                os.remove(image_path)
        except Exception as e:
            print(f"Error adding image: {e}")

        return slide

    def generate_presentation(self, topic, num_slides=5, output_path="generated_presentation.pptx"):
        """Generate a complete PowerPoint presentation"""
        print(f"Generating presentation on: {topic}")

        # Generate content outline
        outline = self.generate_content_outline(topic, num_slides)

        # Create slides based on outline
        for i, slide_data in enumerate(outline):
            title = slide_data["title"]
            content = slide_data["content"]
            slide_type = slide_data["slide_type"]

            print(f"Creating slide {i+1}: {title}")

            if i == 0 or slide_type == "title":
                self.create_title_slide(title, f"Generated by Gemini AI")
            elif slide_type == "image":
                image_query = self.generate_image_description(content)
                self.create_image_slide(title, content, image_query)
            else:
                include_image = (i % 3 == 0)  # Add image every 3rd slide
                self.create_content_slide(title, content, include_image)

        # Save presentation
        self.presentation.save(output_path)
        print(f"Presentation saved as: {output_path}")

        return output_path

print("✅ PPTGenerator class defined successfully!")

api_key = ""
if api_key:
    print("✅ Gemini API key is configured")
else:
    print("❌ Please set your GEMINI_API_KEY in the .env file or uncomment the line above")


    # Initialize the generator
try:
    generator = PPTGenerator(api_key=api_key)
    print("✅ PPT Generator initialized successfully!")
except ValueError as e:
    print(f"❌ Error: {e}")
    print("Please set your GEMINI_API_KEY first.")


    # Generate a presentation
topic = "HITLER"  # Change this to your desired topic
num_slides = 4  # Change this to your desired number of slides

try:
    output_file = generator.generate_presentation(topic, num_slides, "hitler.pptx")
except Exception as e:
    print(e)